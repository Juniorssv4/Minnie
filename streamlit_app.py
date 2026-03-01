import streamlit as st
import time
import google.generativeai as genai
import requests
from io import BytesIO
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential, RetryError
import uuid

# GEMINI CONFIG
try:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
except:
    st.error("Add GEMINI_API_KEY in Secrets")
    st.stop()

# Primary and fallback models
PRIMARY_MODEL = "gemini-2.5-flash"
FALLBACK_MODEL = "gemini-1.5-flash"

if "current_model" not in st.session_state:
    st.session_state.current_model = PRIMARY_MODEL

model = genai.GenerativeModel(st.session_state.current_model)

# Backoff for rate limits
@retry(
    stop=stop_after_attempt(6),
    wait=wait_exponential(multiplier=1, min=4, max=60)
)
def safe_generate_content(prompt):
    return model.generate_content(prompt)

# Glossary from repo file – FORCE FRESH LOAD EVERY TIME, NO CACHE
try:
    # Super strong cache-bust to break GitHub CDN and browser cache
    cache_bust = f"{int(time.time() * 1000)}_{str(uuid.uuid4())[:8]}_{str(hash(str(time.time())))[-6:]}"
    raw_url = f"https://raw.githubusercontent.com/Juniorssv4/minnie-final/main/glossary.txt?cachebust={cache_bust}"
    response = requests.get(raw_url, timeout=10)
    response.raise_for_status()
    lines = response.text.splitlines()
    glossary_dict = {}
    for line in lines:
        line = line.strip()
        if line and ":" in line:
            parts = line.split(":", 1)
            eng = parts[0].strip().lower()
            lao = parts[1].strip() if len(parts) > 1 else ""
            glossary_dict[eng] = lao
    glossary = glossary_dict
except Exception as e:
    glossary = {}
    st.error(f"Glossary load failed: {str(e)}")

def get_glossary_prompt():
    if glossary:
        terms = "\n".join([f"• {e.capitalize()} → {l}" for e, l in glossary.items()])
        return f"Use EXACTLY these terms:\n{terms}\n"
    return ""

def translate_text(text, direction):
    if not text.strip():
        return ""
    target = "Lao" if direction == "English → Lao" else "English"
    prompt = f"""{get_glossary_prompt()}Translate ONLY the text to {target}.
Return ONLY the translation.
Text: {text}"""
    try:
        response = safe_generate_content(prompt)
        return response.text.strip()
    except RetryError as e:
        if "429" in str(e.last_attempt.exception()) or "quota" in str(e.last_attempt.exception()).lower():
            if st.session_state.current_model == PRIMARY_MODEL:
                st.session_state.current_model = FALLBACK_MODEL
                st.info("Rate limit on gemini-2.5-flash — switched to gemini-1.5-flash.")
                global model
                model = genai.GenerativeModel(FALLBACK_MODEL)
                response = model.generate_content(prompt)
                return response.text.strip()
        st.error("Timed out after retries — try again in 5 minutes.")
        return "[Failed — try later]"
    except Exception as e:
        st.error(f"API error: {str(e)}")
        return "[Failed — try again]"

# UI
st.set_page_config(
    page_title="Minnie",
    page_icon="🐶",  # Cute dog emoji as icon
    layout="centered"
)

st.title("🐶 Minnie — Meena's Translator")

direction = st.radio("Direction", ["English → Lao", "Lao → English"], horizontal=True)

tab1, tab2 = st.tabs(["Translate Text", "Translate File"])

with tab1:
    text = st.text_area("Enter text to translate", height=200)
    if st.button("Translate Text", type="primary"):
        with st.spinner("Translating..."):
            result = translate_text(text, direction)
            st.success("Translation:")
            
            # Show translated text
            st.markdown("**Translated text:**")
            st.code(result, language=None)
            
            # Copy button with JS + green success feedback
            copy_js = f"""
                <button onclick="navigator.clipboard.writeText(`{result.replace('`', '\\`').replace('"', '\\"')}`).then(() => {{
                    document.getElementById('copy-success').style.display = 'block';
                    setTimeout(() => {{ document.getElementById('copy-success').style.display = 'none'; }}, 3000);
                }})">📋 Copy to Clipboard</button>
                <p id="copy-success" style="color:green; display:none; margin-top:8px; font-weight:bold;">✅ Copied!</p>
            """
            st.components.v1.html(copy_js, height=60)

with tab2:
    uploaded_file = st.file_uploader("Upload DOCX • XLSX • PPTX (max 50MB)", type=["docx", "xlsx", "pptx"])
    if uploaded_file:
        MAX_SIZE_MB = 50
        if uploaded_file.size > MAX_SIZE_MB * 1024 * 1024:
            st.error(f"File too large! Max allowed size is {MAX_SIZE_MB}MB. Your file is {uploaded_file.size / (1024*1024):.1f}MB.")
        elif st.button("Translate File", type="primary"):
            with st.spinner("Translating file..."):
                file_bytes = uploaded_file.read()
                file_name = uploaded_file.name
                ext = file_name.rsplit(".", 1)[-1].lower()
                output = BytesIO()
                total_elements = 0
                elements_list = []
                if ext == "docx":
                    doc = Document(BytesIO(file_bytes))
                    for p in doc.paragraphs:
                        if p.text.strip():
                            total_elements += 1
                            elements_list.append(("para", p))
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for p in cell.paragraphs:
                                    if p.text.strip():
                                        total_elements += 1
                                        elements_list.append(("para", p))
                elif ext == "xlsx":
                    wb = load_workbook(BytesIO(file_bytes))
                    for ws in wb.worksheets:
                        for row in ws.iter_rows():
                            for cell in row:
                                if isinstance(cell.value, str) and cell.value.strip():
                                    total_elements += 1
                                    elements_list.append(("cell", cell))
                elif ext == "pptx":
                    prs = Presentation(BytesIO(file_bytes))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for p in shape.text_frame.paragraphs:
                                    if p.text.strip():
                                        total_elements += 1
                                        elements_list.append(("para", p))
                if total_elements == 0:
                    st.warning("No text found in file.")
                    st.stop()
                progress_bar = st.progress(0)
                status_text = st.empty()
                translated_count = 0
                for element_type, element in elements_list:
                    status_text.text(f"Translating... {translated_count}/{total_elements}")
                    if element_type == "para":
                        translated = translate_text(element.text, direction)
                        element.text = translated
                    elif element_type == "cell":
                        translated = translate_text(element.value, direction)
                        element.value = translated
                    translated_count += 1
                    progress_bar.progress(translated_count / total_elements)
                status_text.text("Saving file...")
                if ext == "docx":
                    doc.save(output)
                elif ext == "xlsx":
                    wb.save(output)
                elif ext == "pptx":
                    prs.save(output)
                output.seek(0)
                filename = f"TRANSLATED_{file_name}"
                mime_type = "application/octet-stream"
                st.success("Translation complete!")
                st.info("Click the big button below to download your translated file. Your browser may block auto-downloads — this button always works!")
                st.download_button(
                    label="📥 DOWNLOAD TRANSLATED FILE NOW",
                    data=output,
                    file_name=filename,
                    mime=mime_type,
                    type="primary",
                    use_container_width=True,
                    key="download_btn_" + str(time.time()),
                    help="Click here to save the translated file to your device"
                )
                st.caption("Tip: If nothing happens, refresh the page or try in another browser (Chrome works best).")

# Teach term (manual in GitHub)
with st.expander("➕ Teach Minnie a new term (edit glossary.txt in GitHub)"):
    st.info("To add term: Edit glossary.txt in repo → add line 'english:lao' → save → click the red reload button below or refresh page.")
    st.code("Example:\nSamir:ສະຫມີຣ\nhello:ສະບາຍດີ")

# Big red manual reload button – click this RIGHT AFTER you commit changes to glossary.txt
st.markdown("---")
st.markdown("<h3 style='color:red;'>If you just edited glossary.txt on GitHub, click this button now:</h3>", unsafe_allow_html=True)
if st.button("🔴 RELOAD GLOSSARY FROM GITHUB (click after editing & committing)", type="primary", use_container_width=True):
    st.rerun()

st.caption(f"Active glossary: {len(glossary)} terms • Model: {st.session_state.current_model}")
